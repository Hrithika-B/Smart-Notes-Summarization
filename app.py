import streamlit as st
from summarizer import TextSummarizer
import PyPDF2, docx
from pptx import Presentation

st.set_page_config(page_title="AI Text Summarizer", layout="centered")

# ---------------- SESSION STATE ----------------
if "logged_in" not in st.session_state:
    st.session_state["logged_in"] = False
if "username" not in st.session_state:
    st.session_state["username"] = ""
if "uploaded_text" not in st.session_state:
    st.session_state["uploaded_text"] = ""

# ---------------- LOGIN ----------------
USERS = {"admin": "12345"}

def do_login(username_input, password_input):
    if username_input in USERS and password_input == USERS[username_input]:
        st.session_state["logged_in"] = True
        st.session_state["username"] = username_input
    else:
        st.session_state["logged_in"] = False
        st.error("Invalid username or password ❌")

def login_page():
    st.title("🔐 Login Page")
    st.markdown("Please enter your credentials to access the AI Text Summarizer Dashboard.")
    username_input = st.text_input("Username", key="login_username")
    password_input = st.text_input("Password", type="password", key="login_password")
    st.button("Login", on_click=do_login, args=(username_input, password_input))

# ---------------- FILE READING ----------------
def read_txt(file): return file.read().decode("utf-8")
def read_pdf(file):
    pdf_reader = PyPDF2.PdfReader(file)
    return "\n".join([page.extract_text() for page in pdf_reader.pages if page.extract_text()])
def read_docx(file):
    doc = docx.Document(file)
    return "\n".join([para.text for para in doc.paragraphs])
def read_pptx(file):
    prs = Presentation(file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

# ---------------- SUMMARIZER ----------------
@st.cache_resource
def load_summarizer():
    return TextSummarizer(model_name="facebook/bart-large-cnn")

def chunk_text(text, chunk_size=800):
    words = text.split()
    return [" ".join(words[i:i+chunk_size]) for i in range(0, len(words), chunk_size)]

def summarize_large_text(summarizer, text, max_len=300, min_len=50):
    chunks = chunk_text(text, chunk_size=800)
    summaries = []
    for chunk in chunks:
        summary = summarizer.summarize_text(chunk, max_length=max_len, min_length=min_len)
        summaries.append(summary)
    # Optional: summarize all chunk summaries again for concise final summary
    combined_text = " ".join(summaries)
    final_summary = summarizer.summarize_text(combined_text, max_length=max_len, min_length=min_len)
    return final_summary

# ---------------- DASHBOARD ----------------
def dashboard():
    st.title("💡 AI Text Summarizer Dashboard")
    st.markdown(f"Welcome, **{st.session_state['username']}**!")

    summarizer_instance = load_summarizer()
    if not summarizer_instance.summarizer:
        st.error("Failed to load the summarization model.")
        return

    uploaded_file = st.file_uploader(
        "Upload your notes (.txt, .pdf, .docx, .pptx)", 
        type=["txt", "pdf", "docx", "pptx"]
    )

    if uploaded_file is not None:
        file_type = uploaded_file.name.split(".")[-1].lower()
        if file_type == "txt": text = read_txt(uploaded_file)
        elif file_type == "pdf": text = read_pdf(uploaded_file)
        elif file_type == "docx": text = read_docx(uploaded_file)
        elif file_type == "pptx": text = read_pptx(uploaded_file)
        else:
            st.error("Unsupported file type!")
            return

        st.session_state["uploaded_text"] = text
        st.success(f"{uploaded_file.name} uploaded successfully! ✅")
        st.text_area("Preview of uploaded text:", value=text[:1000], height=200)

    if st.session_state["uploaded_text"]:
        max_len = st.slider("Maximum summary length", 50, 500, 300)
        min_len = st.slider("Minimum summary length", 10, 200, 50)
        do_sample = st.checkbox("Enable creative sampling (may vary results)", value=False)

        if st.button("Summarize Notes"):
            with st.spinner("Generating summary... This may take a few seconds for large files."):
                summary = summarize_large_text(
                    summarizer_instance,
                    st.session_state["uploaded_text"],
                    max_len=max_len,
                    min_len=min_len
                )
                st.subheader("📄 Summary:")
                st.write(summary)

    st.markdown("""
    <hr>
    <p style='text-align: center;'>Powered by Hugging Face Transformers and Streamlit</p>
    <p style='text-align: center;'>Developed by Hrithika & Ramya K M</p>
    """, unsafe_allow_html=True)

# ---------------- MAIN ----------------
if not st.session_state["logged_in"]:
    login_page()
else:
    dashboard()
